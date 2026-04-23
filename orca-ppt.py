import os

try:
    from pptx import Presentation
    from pptx.util import Pt
except Exception as e:
    raise RuntimeError("python-pptx is required to run this script. Install it with: pip install python-pptx")

slides = [
    ("ORCA", "Object Reconstruction via Computational AI\nAdvanced Project Review"),
    ("Introduction", "AI-powered 3D reconstruction system\nTransforms multi-view images into 3D models\nCombines Photogrammetry + Generative AI"),
    ("Problem Statement", "Existing tools fail on incomplete objects\nRequire expert-level configuration\nNo automation for full pipeline\nNo free scalable AI solutions"),
    ("System Architecture", "Client Interface (React + Three.js)\nFastAPI Backend\nJob Management System\nReconstruction Pipeline\nCloud Storage (Supabase)"),
    ("Method 1 - Photogrammetry", "COLMAP SfM (Sparse Reconstruction)\nDense Stereo Reconstruction\nPoisson Mesh Generation\nMesh Cleaning & Optimization"),
    ("Method 2 - Hybrid AI", "COLMAP Pose Estimation\nNeRF Neural Reconstruction\nHole Detection\nLocal AI Completion\nMesh Fusion & Texture Baking"),
    ("Key Innovation", "No reliance on paid APIs\nFully local AI pipeline\nHybrid reconstruction approach\nDomain-focused optimization (objects like shells)"),
    ("Challenges Faced", "Lack of free high-quality 3D AI APIs\nTrial APIs insufficient\nGPU memory limitations\nHandling incomplete geometry"),
    ("Results", "Generated high-quality meshes\nQuality metrics evaluation\nSuccessful reconstruction pipeline\nSupports multiple outputs (.ply, .obj)"),
    ("Limitations", "Requires GPU for processing\nNot real-time\nLimited domain currently\nOccasional non-watertight meshes"),
    ("Future Scope", "Full NeRF pipeline integration\nAdvanced AI completion models\nMobile capture system\nCloud GPU deployment"),
    ("Conclusion", "Scalable AI-assisted 3D reconstruction system\nBridges gap between photogrammetry and AI\nStrong foundation for future research")
]

prs = Presentation()

for i, (title, content) in enumerate(slides):
    if i == 0:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
    else:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        first = True
        for line in content.split("\n"):
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(20)

output_dir = os.path.join(os.getcwd(), "output")
os.makedirs(output_dir, exist_ok=True)

file_path = os.path.join(output_dir, "ORCA_Advanced_Presentation.pptx")
prs.save(file_path)

assert os.path.isdir(output_dir)
assert os.path.exists(file_path)
assert os.path.getsize(file_path) > 0

file_path
